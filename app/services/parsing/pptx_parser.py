"""PowerPoint (.pptx) parser implementing ParserInterface.

Uses python-pptx (lazily imported) to extract per-slide text, tables, and
speaker notes. A missing python-pptx dependency or unexpected error degrades
gracefully into a ``ParsedDocument`` with a warning rather than raising.
"""

from __future__ import annotations

import time
from pathlib import Path

from app.services.parsing.base import PageContent, ParsedDocument, ParserInterface


class PptxParser(ParserInterface):
    """Parse ``.pptx`` presentations: one page per slide."""

    supported_extensions = [".pptx"]

    async def parse(self, file_path: str) -> ParsedDocument:
        start = time.monotonic()
        path = Path(file_path)

        try:
            from pptx import Presentation
        except ImportError:
            elapsed = int((time.monotonic() - start) * 1000)
            return ParsedDocument(
                filename=path.name,
                content_type="pptx",
                full_text="",
                pages=[],
                tables=[],
                metadata={},
                page_count=0,
                processing_time_ms=elapsed,
                warnings=["python-pptx not installed; .pptx content not extracted"],
            )

        warnings: list[str] = []
        try:
            prs = Presentation(file_path)

            pages: list[PageContent] = []
            all_tables: list[dict] = []
            slide_texts: list[str] = []

            for slide_num, slide in enumerate(prs.slides, 1):
                parts: list[str] = []
                page_tables: list[dict] = []

                for shape in slide.shapes:
                    if getattr(shape, "has_table", False):
                        rows = [
                            [cell.text.strip() for cell in row.cells]
                            for row in shape.table.rows
                        ]
                        if rows:
                            headers = rows[0]
                            data = rows[1:]
                            table_dict = {
                                "slide": slide_num,
                                "headers": headers,
                                "data": data,
                                "rows": len(rows),
                                "cols": len(headers),
                            }
                            page_tables.append(table_dict)
                            all_tables.append(table_dict)
                            for row in rows:
                                parts.append(" | ".join(row))
                    elif hasattr(shape, "text") and shape.text:
                        parts.append(shape.text)

                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        parts.append(f"[Notes: {notes}]")

                slide_text = "\n".join(parts)
                slide_texts.append(slide_text)
                pages.append(
                    PageContent(
                        page_number=slide_num,
                        text=slide_text,
                        tables=page_tables,
                    )
                )

            full_text = "\n\n".join(slide_texts)
            props = prs.core_properties
            metadata = {
                "title": props.title or None,
                "author": props.author or None,
                "slide_count": len(prs.slides),
            }

            elapsed = int((time.monotonic() - start) * 1000)
            return ParsedDocument(
                filename=path.name,
                content_type="pptx",
                full_text=full_text,
                pages=pages,
                tables=all_tables,
                metadata=metadata,
                page_count=len(pages),
                processing_time_ms=elapsed,
                warnings=warnings,
            )

        except Exception as exc:  # noqa: BLE001 -- never crash ingestion.
            elapsed = int((time.monotonic() - start) * 1000)
            return ParsedDocument(
                filename=path.name,
                content_type="pptx",
                full_text="",
                pages=[],
                tables=[],
                metadata={},
                page_count=0,
                processing_time_ms=elapsed,
                warnings=[f"Parse error: {exc}"],
            )
