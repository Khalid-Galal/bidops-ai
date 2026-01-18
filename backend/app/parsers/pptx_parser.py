"""PowerPoint presentation parser using python-pptx."""

import time
from pathlib import Path

from app.parsers.base import BaseParser, ParsedContent, ParserRegistry


@ParserRegistry.register
class PptxParser(BaseParser):
    """Parser for Microsoft PowerPoint documents (.pptx)."""

    supported_extensions = [".pptx", ".ppt"]
    supported_mimetypes = [
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint",
    ]

    async def parse(self, file_path: str) -> ParsedContent:
        """Parse a PowerPoint presentation.

        Args:
            file_path: Path to PPTX file

        Returns:
            ParsedContent with extracted text and metadata
        """
        from pptx import Presentation
        from pptx.util import Inches

        self.validate_file(file_path)
        start_time = time.time()

        try:
            prs = Presentation(file_path)

            slides_text = []
            tables = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = [f"\n=== Slide {slide_num} ===\n"]

                for shape in slide.shapes:
                    # Extract text from shapes
                    if hasattr(shape, "text") and shape.text:
                        slide_content.append(shape.text)

                    # Extract tables
                    if shape.has_table:
                        table_data = []
                        for row in shape.table.rows:
                            row_data = [
                                cell.text.strip() for cell in row.cells
                            ]
                            table_data.append(row_data)
                            slide_content.append(" | ".join(row_data))

                        tables.append({
                            "slide": slide_num,
                            "data": table_data,
                        })

                # Get notes if available
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        slide_content.append(f"\n[Notes: {notes}]")

                slides_text.append("\n".join(slide_content))

            full_text = "\n\n".join(slides_text)

            # Extract metadata
            metadata = await self.extract_metadata(file_path)
            metadata["slide_count"] = len(prs.slides)

            processing_time = int((time.time() - start_time) * 1000)

            return ParsedContent(
                text=full_text,
                metadata=metadata,
                pages=slides_text,
                tables=tables if tables else None,
                page_count=len(prs.slides),
                processing_time_ms=processing_time,
            )

        except Exception as e:
            raise Exception(f"Failed to parse PowerPoint: {str(e)}") from e

    async def extract_metadata(self, file_path: str) -> dict:
        """Extract metadata from PowerPoint file.

        Args:
            file_path: Path to PPTX file

        Returns:
            Dictionary of metadata
        """
        from pptx import Presentation

        try:
            prs = Presentation(file_path)
            props = prs.core_properties

            return {
                "title": props.title or None,
                "author": props.author or None,
                "subject": props.subject or None,
                "keywords": props.keywords or None,
                "created": props.created.isoformat() if props.created else None,
                "modified": props.modified.isoformat() if props.modified else None,
                "last_modified_by": props.last_modified_by or None,
                "category": props.category or None,
                "file_size": Path(file_path).stat().st_size,
            }
        except Exception:
            return {
                "file_size": Path(file_path).stat().st_size,
            }
