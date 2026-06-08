async def test_pptx_parser_extracts_slide_text(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches
    from app.services.parsing.pptx_parser import PptxParser

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    box.text_frame.text = "Project Kickoff - New Cairo Medical Center"
    f = tmp_path / "deck.pptx"
    prs.save(str(f))

    parsed = await PptxParser().parse(str(f))
    assert parsed.content_type == "pptx"
    assert "New Cairo Medical Center" in parsed.full_text
    assert parsed.page_count == 1
    assert parsed.pages[0].page_number == 1
